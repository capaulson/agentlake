"""PowerPoint (.pptx) file adapter using python-pptx."""

from __future__ import annotations

from pathlib import Path

from agentlake.adapters.base import ExtractedContent, StructureMarker, TextBlock


class PptxAdapter:
    """Extracts text from PowerPoint presentations.

    Text is extracted from each slide's shapes with source locator
    ``slide:{n}`` (1-indexed).
    """

    supported_extensions: list[str] = [".pptx"]
    supported_mimetypes: list[str] = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]

    def can_handle(self, filename: str, content_type: str) -> bool:
        return (
            Path(filename).suffix.lower() in self.supported_extensions
            or content_type in self.supported_mimetypes
        )

    def extract(self, file_bytes: bytes, filename: str) -> ExtractedContent:
        """Extract text from each slide.

        Args:
            file_bytes: Raw PPTX bytes.
            filename: Original filename.

        Returns:
            Extracted content with per-slide text blocks.
        """
        import io

        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))

        metadata: dict = {
            "filename": filename,
            "slide_count": len(prs.slides),
        }

        blocks: list[TextBlock] = []
        structure: list[StructureMarker] = []
        position = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            structure.append(
                StructureMarker(
                    marker_type="section_start",
                    position=position,
                    metadata={"slide": slide_idx},
                )
            )

            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        row_text = " | ".join(cells)
                        if row_text.strip():
                            slide_texts.append(row_text)

            if slide_texts:
                blocks.append(
                    TextBlock(
                        content="\n".join(slide_texts),
                        block_type="paragraph",
                        position=position,
                        source_locator=f"slide:{slide_idx}",
                        metadata={"slide": slide_idx},
                    )
                )
                position += 1

        return ExtractedContent(
            text_blocks=blocks,
            metadata=metadata,
            structure=structure,
        )
